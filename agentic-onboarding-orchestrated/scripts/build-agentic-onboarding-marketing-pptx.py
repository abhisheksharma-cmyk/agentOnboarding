from __future__ import annotations

from datetime import datetime
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def latest_merged_docx(repo: Path) -> Path:
    matches = sorted(repo.glob("AGENTIC_ONBOARDING_MERGED_BRIEF*.docx"), key=lambda p: p.stat().st_mtime)
    if not matches:
        raise FileNotFoundError("No AGENTIC_ONBOARDING_MERGED_BRIEF*.docx found.")
    return matches[-1]


def read_docx_paragraphs(docx_path: Path) -> list[str]:
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(docx_path) as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)
    lines: list[str] = []
    for p in root.findall(".//w:p", ns):
        text = "".join((n.text or "") for n in p.findall(".//w:t", ns)).strip()
        if text:
            lines.append(text)
    return lines


def section_between(lines: list[str], start_prefix: str, end_prefixes: tuple[str, ...]) -> list[str]:
    start_idx = next((i for i, s in enumerate(lines) if s.startswith(start_prefix)), -1)
    if start_idx == -1:
        return []
    out: list[str] = []
    for s in lines[start_idx + 1 :]:
        if any(s.startswith(ep) for ep in end_prefixes):
            break
        out.append(s)
    return out


def fit_into_box(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int]:
    r_img = img_w / img_h
    r_box = box_w / box_h
    if r_img > r_box:
        w = box_w
        h = int(w / r_img)
    else:
        h = box_h
        w = int(h * r_img)
    return w, h


def set_slide_bg(slide, rgb=(246, 250, 255)) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*rgb)


def add_header_band(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(17, 38, 74)
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8.8), Inches(0.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.66), Inches(10.0), Inches(0.24))
        sf = sb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Arial"
        sp.font.size = Pt(11)
        sp.font.color.rgb = RGBColor(209, 213, 219)


def add_footer(slide, text: str) -> None:
    fb = slide.shapes.add_textbox(Inches(0.45), Inches(7.18), Inches(12.4), Inches(0.2))
    tf = fb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation, generated: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (244, 248, 255))

    hero = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(240, 246, 255)
    hero.line.fill.background()

    stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(12, 30, 61)
    stripe.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(11.8), Inches(0.5))
    tf1 = t1.text_frame
    tf1.clear()
    p1 = tf1.paragraphs[0]
    p1.text = "Verinite Agentic Onboarding"
    p1.font.name = "Arial"
    p1.font.bold = True
    p1.font.size = Pt(34)
    p1.font.color.rgb = RGBColor(255, 255, 255)

    t2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(12.0), Inches(2.0))
    tf2 = t2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "Accelerating compliant customer onboarding through policy-governed Agentic AI."
    p2.font.name = "Arial"
    p2.font.bold = True
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(17, 38, 74)

    p3 = tf2.add_paragraph()
    p3.text = "For enterprise banks, NBFCs, fintechs, and regulated digital businesses."
    p3.font.name = "Arial"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(51, 65, 85)
    p3.space_before = Pt(12)

    add_footer(slide, f"Confidential | Verinite | Generated {generated}")


def add_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(slide, "Agenda", "Client-focused transformation narrative")

    items = [
        "Why onboarding transformation is urgent",
        "Current challenges in customer onboarding operations",
        "Agentic Onboarding solution overview",
        "Business impact and value realization",
        "Architecture, control, and governance model",
        "Implementation roadmap with Verinite services",
    ]

    panel = slide.shapes.add_shape(1, Inches(0.8), Inches(1.35), Inches(11.7), Inches(5.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(203, 213, 225)

    tf = panel.text_frame
    tf.clear()
    for idx, item in enumerate(items, 1):
        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
        p.text = f"{idx:02d}. {item}"
        p.font.name = "Arial"
        p.font.size = Pt(20 if idx == 1 else 18)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_before = Pt(8 if idx > 1 else 0)

    add_footer(slide, "Verinite Agentic Onboarding | Agenda")


def add_challenges_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(
        slide,
        "Understanding the Current Challenge",
        "Operational pain points in traditional onboarding models",
    )

    left = slide.shapes.add_shape(1, Inches(0.65), Inches(1.35), Inches(6.1), Inches(5.75))
    right = slide.shapes.add_shape(1, Inches(6.95), Inches(1.35), Inches(5.75), Inches(5.75))
    for s in (left, right):
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(255, 255, 255)
        s.line.color.rgb = RGBColor(203, 213, 225)

    tf_l = left.text_frame
    tf_l.clear()
    h = tf_l.paragraphs[0]
    h.text = "Challenge"
    h.font.name = "Arial"
    h.font.bold = True
    h.font.size = Pt(22)
    h.font.color.rgb = RGBColor(15, 23, 42)
    for t in [
        "Manual and repetitive verification activities",
        "Delayed turnaround for known onboarding scenarios",
        "Scattered operational knowledge across teams and channels",
        "High dependence on specialist intervention for basic decisions",
    ]:
        p = tf_l.add_paragraph()
        p.text = t
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(7)

    tf_r = right.text_frame
    tf_r.clear()
    h2 = tf_r.paragraphs[0]
    h2.text = "Business Impact"
    h2.font.name = "Arial"
    h2.font.bold = True
    h2.font.size = Pt(22)
    h2.font.color.rgb = RGBColor(15, 23, 42)
    for t in [
        "Higher onboarding cost per application",
        "Inconsistent processing quality and customer experience",
        "Compliance risk due to non-standardized execution",
        "Reduced growth velocity because of slower onboarding throughput",
    ]:
        p = tf_r.add_paragraph()
        p.text = t
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(7)

    add_footer(slide, "Verinite Agentic Onboarding | Challenge Landscape")


def add_solution_overview_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(
        slide,
        "Reimagining Onboarding with Agentic AI",
        "A policy-governed, event-driven decision system",
    )

    intro = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(12.0), Inches(0.55))
    tf = intro.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "Agentic Onboarding combines specialized AI agents, deterministic workflow control, and centralized policy governance "
        "to deliver faster, safer, and more explainable customer onboarding decisions."
    )
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 65, 85)

    cards = [
        "1. Frontend Data and Document Intake",
        "2. Orchestrator + Trace-Aware State Machine",
        "3. KYC and Address Verification Agents",
        "4. AML and Credit Decision Agents",
        "5. Risk Synthesis + Decision Gateway",
        "6. Full Audit Trail and Explainable Output",
    ]
    x0, y0 = 0.75, 2.0
    w, h = 3.95, 1.65
    gapx, gapy = 0.28, 0.35
    for i, txt in enumerate(cards):
        col = i % 3
        row = i // 3
        x = Inches(x0 + col * (w + gapx))
        y = Inches(y0 + row * (h + gapy))
        box = slide.shapes.add_shape(1, x, y, Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(191, 219, 254)
        btf = box.text_frame
        btf.clear()
        bp = btf.paragraphs[0]
        bp.text = txt
        bp.font.name = "Arial"
        bp.font.bold = True
        bp.font.size = Pt(14)
        bp.font.color.rgb = RGBColor(29, 78, 216)

    add_footer(slide, "Verinite Agentic Onboarding | Solution Overview")


def add_value_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(
        slide,
        "Business Value and Outcome Metrics",
        "Indicative targets based on automated onboarding operating models",
    )

    metrics = [
        ("30-50%", "Faster application decision turnaround"),
        ("40-60%", "Reduction in repetitive manual processing"),
        ("100%", "Traceable audit chain for every decision"),
        ("High", "Regulatory confidence with policy-governed controls"),
    ]

    x0 = 0.9
    for i, (num, desc) in enumerate(metrics):
        x = Inches(x0 + i * 3.05)
        card = slide.shapes.add_shape(1, x, Inches(2.0), Inches(2.75), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(191, 219, 254)

        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = num
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(34)
        p1.font.color.rgb = RGBColor(17, 38, 74)
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Arial"
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(51, 65, 85)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)

    disclaimer = slide.shapes.add_textbox(Inches(0.75), Inches(6.6), Inches(12.0), Inches(0.3))
    dtf = disclaimer.text_frame
    dtf.clear()
    dp = dtf.paragraphs[0]
    dp.text = "Note: Metrics are indicative and refined during discovery and pilot baselining."
    dp.font.name = "Arial"
    dp.font.size = Pt(10)
    dp.font.color.rgb = RGBColor(100, 116, 139)

    add_footer(slide, "Verinite Agentic Onboarding | Business Outcomes")


def add_image_slide(prs: Presentation, title: str, subtitle: str, image_path: Path, footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(slide, title, subtitle)

    left = Inches(0.55)
    top = Inches(1.12)
    box_w = prs.slide_width - Inches(1.1)
    box_h = prs.slide_height - Inches(1.65)

    with Image.open(image_path) as img:
        iw, ih = img.size
    w, h = fit_into_box(iw, ih, int(box_w), int(box_h))
    x = int(left + (box_w - w) / 2)
    y = int(top + (box_h - h) / 2)
    slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)

    add_footer(slide, footer)


def add_roadmap_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_band(
        slide,
        "Engagement Roadmap",
        "A phased approach for rapid value realization with controlled risk",
    )

    phases = [
        ("Phase 1", "Discover", "Process assessment, compliance baseline, integration mapping"),
        ("Phase 2", "Pilot", "Limited-scope rollout with trace, policy, and outcome validation"),
        ("Phase 3", "Scale", "Progressive expansion across segments and products"),
        ("Phase 4", "Optimize", "Continuous tuning, model-policy refinement, KPI governance"),
    ]

    y = Inches(2.0)
    for i, (ph, name, desc) in enumerate(phases):
        x = Inches(0.8 + i * 3.1)
        box = slide.shapes.add_shape(1, x, y, Inches(2.75), Inches(3.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(191, 219, 254)
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = ph
        p0.font.name = "Arial"
        p0.font.bold = True
        p0.font.size = Pt(13)
        p0.font.color.rgb = RGBColor(29, 78, 216)
        p1 = tf.add_paragraph()
        p1.text = name
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = RGBColor(15, 23, 42)
        p1.space_before = Pt(3)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Arial"
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(51, 65, 85)
        p2.space_before = Pt(8)

    add_footer(slide, "Verinite Agentic Onboarding | Engagement Roadmap")


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (241, 245, 249))

    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(17, 38, 74)
    band.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.8), Inches(2.0))
    tf = t.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Thank You"
    p1.font.name = "Arial"
    p1.font.bold = True
    p1.font.size = Pt(54)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Let us co-design your next-generation onboarding experience."
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(191, 219, 254)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    c = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.4))
    ctf = c.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = "Verinite | IT Services and Digital Transformation"
    cp.font.name = "Arial"
    cp.font.size = Pt(12)
    cp.font.color.rgb = RGBColor(209, 213, 219)
    cp.alignment = PP_ALIGN.CENTER


def main() -> None:
    repo = Path.cwd()
    docx_path = latest_merged_docx(repo)
    lines = read_docx_paragraphs(docx_path)

    strategic = section_between(lines, "1) System Purpose and Strategic Intent", ("2) Consolidated Architecture View",))
    features = section_between(lines, "6) Feature Highlights", ("7) Implementation Anchors",))
    _ = strategic, features  # parsed for traceability; slide text is curated from merged brief narrative.

    diag_dir = repo / "documents" / "diagrams"
    diag_topology = diag_dir / "architecture-topology.png"
    diag_sequence = diag_dir / "orchestration-sequence.png"
    diag_yaml = diag_dir / "yaml-pluggability-control-surface.png"
    for p in (diag_topology, diag_sequence, diag_yaml):
        if not p.exists():
            raise FileNotFoundError(f"Missing diagram image: {p}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    generated = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(prs, generated)
    add_agenda_slide(prs)
    add_challenges_slide(prs)
    add_solution_overview_slide(prs)
    add_value_slide(prs)
    add_image_slide(
        prs,
        "Architecture Overview",
        "Event-driven orchestration with centralized decision governance",
        diag_topology,
        "Verinite Agentic Onboarding | Architecture Overview",
    )
    add_image_slide(
        prs,
        "Workflow and Decision Progression",
        "Deterministic stage transitions with policy guardrails",
        diag_sequence,
        "Verinite Agentic Onboarding | Workflow and Decisioning",
    )
    add_image_slide(
        prs,
        "Configuration and Pluggability Control Surface",
        "YAML-driven slot versioning and runtime operational controls",
        diag_yaml,
        "Verinite Agentic Onboarding | YAML Pluggability",
    )
    add_roadmap_slide(prs)
    add_closing_slide(prs)

    out_name = f"AGENTIC_ONBOARDING_MARKETING_DECK_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    out_path = repo / out_name
    prs.save(out_path)
    print(out_path)


if __name__ == "__main__":
    main()
