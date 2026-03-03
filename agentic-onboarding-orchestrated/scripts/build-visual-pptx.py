from __future__ import annotations

from datetime import datetime
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def fit_into_box(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int]:
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        w = box_w
        h = int(w / img_ratio)
    else:
        h = box_h
        w = int(h * img_ratio)
    return w, h


def read_docx_paragraphs(docx_path: Path) -> list[str]:
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(docx_path) as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)
    paragraphs: list[str] = []
    for p in root.findall(".//w:p", ns):
        text = "".join((node.text or "") for node in p.findall(".//w:t", ns)).strip()
        if text:
            paragraphs.append(text)
    return paragraphs


def latest_visual_docx(repo: Path) -> Path:
    matches = sorted(
        repo.glob("AGENTIC_ONBOARDING_VISUAL_BRIEF*.docx"),
        key=lambda p: p.stat().st_mtime,
    )
    if not matches:
        raise FileNotFoundError("No AGENTIC_ONBOARDING_VISUAL_BRIEF*.docx found.")
    return matches[-1]


def section_between(lines: list[str], start_prefix: str, end_prefixes: tuple[str, ...]) -> list[str]:
    start_idx = next((i for i, s in enumerate(lines) if s.startswith(start_prefix)), -1)
    if start_idx == -1:
        return []
    out: list[str] = []
    for s in lines[start_idx + 1 :]:
        if any(s.startswith(prefix) for prefix in end_prefixes):
            break
        out.append(s)
    return out


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246, 250, 255)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.name = "Segoe UI"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)

    sub_tf = subtitle_shape.text_frame
    sub_tf.paragraphs[0].font.name = "Segoe UI"
    sub_tf.paragraphs[0].font.size = Pt(18)
    sub_tf.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)


def add_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 255)

    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Presentation Scope"
    p.font.name = "Segoe UI"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    panel = slide.shapes.add_shape(
        1, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.7)
    )  # MSO_SHAPE.RECTANGLE == 1
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(203, 213, 225)
    panel.line.width = Pt(1)

    panel_tf = panel.text_frame
    panel_tf.clear()
    panel_tf.word_wrap = True

    heading = panel_tf.paragraphs[0]
    heading.text = "This deck includes"
    heading.font.name = "Segoe UI"
    heading.font.size = Pt(20)
    heading.font.bold = True
    heading.font.color.rgb = RGBColor(15, 23, 42)

    items = [
        "1. Platform topology with centralized policy authority and distributed slot execution.",
        "2. End-to-end orchestration sequence with deterministic transitions and guardrails.",
        "3. YAML pluggability control surface with active-version routing and runtime safety controls.",
        "4. Image-first layout with fitted diagrams and on-slide narrative highlights.",
    ]
    for item in items:
        p = panel_tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        p.font.name = "Segoe UI"
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(51, 65, 85)


def add_text_highlight_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 255)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    panel = slide.shapes.add_shape(1, Inches(0.7), Inches(1.15), Inches(11.9), Inches(5.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(203, 213, 225)
    panel.line.width = Pt(1)

    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    if not bullets:
        bullets = ["No content found in source document section."]
    for i, bullet in enumerate(bullets):
        bp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bp.text = bullet
        bp.level = 0
        bp.space_before = Pt(4)
        bp.space_after = Pt(1)
        bp.font.name = "Segoe UI"
        bp.font.size = Pt(18 if i == 0 else 16)
        bp.font.color.rgb = RGBColor(51, 65, 85)


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    summary: str,
    bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(249, 251, 255)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), slide_w - Inches(0.9), Inches(0.55))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.62), slide_w - Inches(0.9), Inches(0.32))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    sp = subtitle_tf.paragraphs[0]
    sp.text = summary
    sp.font.name = "Segoe UI"
    sp.font.size = Pt(13)
    sp.font.color.rgb = RGBColor(71, 85, 105)
    sp.alignment = PP_ALIGN.LEFT

    # Image-first layout: large full-width image area + compact bottom highlight strip
    left_margin = Inches(0.45)
    right_margin = Inches(0.45)
    top = Inches(1.0)
    bottom_margin = Inches(1.35)
    image_box_w = slide_w - left_margin - right_margin

    box_h = slide_h - top - bottom_margin

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    fitted_w, fitted_h = fit_into_box(img_w, img_h, int(image_box_w), int(box_h))
    left = int(left_margin + (image_box_w - fitted_w) / 2)
    image_top = int(top + (box_h - fitted_h) / 2)

    slide.shapes.add_picture(str(image_path), left, image_top, width=fitted_w, height=fitted_h)

    panel_top = slide_h - Inches(1.23)
    panel = slide.shapes.add_shape(1, left_margin, panel_top, image_box_w, Inches(0.85))  # rectangle
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(203, 213, 225)
    panel.line.width = Pt(1)

    panel_tf = panel.text_frame
    panel_tf.clear()
    panel_tf.word_wrap = True
    panel_tf.margin_left = Inches(0.16)
    panel_tf.margin_right = Inches(0.16)
    panel_tf.margin_top = Inches(0.12)
    panel_tf.margin_bottom = Inches(0.08)

    hp = panel_tf.paragraphs[0]
    hp.text = "Key Highlights"
    hp.font.name = "Segoe UI"
    hp.font.bold = True
    hp.font.size = Pt(13)
    hp.font.color.rgb = RGBColor(15, 23, 42)

    compact = []
    for bullet in bullets[:3]:
        normalized = bullet.lstrip("- ").strip()
        if len(normalized) > 78:
            normalized = normalized[:75].rstrip() + "..."
        compact.append(normalized)
    bp = panel_tf.add_paragraph()
    bp.text = " | ".join(compact) if compact else "No highlights available."
    bp.level = 0
    bp.space_before = Pt(2)
    bp.space_after = Pt(0)
    bp.font.name = "Segoe UI"
    bp.font.size = Pt(11)
    bp.font.color.rgb = RGBColor(51, 65, 85)

    # Optional source line
    source_box = slide.shapes.add_textbox(Inches(0.5), slide_h - Inches(0.26), slide_w - Inches(1.0), Inches(0.2))
    source_tf = source_box.text_frame
    source_tf.clear()
    sp = source_tf.paragraphs[0]
    sp.text = f"Source: {image_path.as_posix()}"
    sp.font.name = "Segoe UI"
    sp.font.size = Pt(9)
    sp.font.color.rgb = RGBColor(100, 116, 139)
    sp.alignment = PP_ALIGN.RIGHT


def main() -> None:
    repo = Path.cwd()
    diagram_dir = repo / "documents" / "diagrams"
    source_doc = latest_visual_docx(repo)
    doc_lines = read_docx_paragraphs(source_doc)

    strategic = section_between(
        doc_lines,
        "1) Strategic Narrative",
        ("2) Visual Diagram A: Platform Topology",),
    )
    summary_a = section_between(
        doc_lines,
        "2) Visual Diagram A: Platform Topology",
        ("3) Visual Diagram B: Orchestration Sequence",),
    )
    summary_b = section_between(
        doc_lines,
        "3) Visual Diagram B: Orchestration Sequence",
        ("4) Visual Diagram C: YAML Control Surface",),
    )
    summary_c = section_between(
        doc_lines,
        "4) Visual Diagram C: YAML Control Surface",
        ("5) Why This Is Pluggable and Configurable (YAML Deep Dive)",),
    )
    yaml_bullets = section_between(
        doc_lines,
        "5) Why This Is Pluggable and Configurable (YAML Deep Dive)",
        ("6) Feature Highlights with Enterprise Positioning",),
    )
    feature_bullets = section_between(
        doc_lines,
        "6) Feature Highlights with Enterprise Positioning",
        ("7) Implementation Anchors in Current Codebase",),
    )
    impl_bullets = section_between(
        doc_lines,
        "7) Implementation Anchors in Current Codebase",
        tuple(),
    )

    subtitle = next(
        (s for s in doc_lines if s.startswith("Purpose:")),
        "Visual architecture deck generated from source document narrative.",
    )

    slides = [
        (
            "Platform Topology",
            diagram_dir / "architecture-topology.png",
            summary_a[0] if summary_a else "Distributed decision mesh with centralized policy authority.",
            [
                *(strategic[:2] if strategic else []),
                "Decision Gateway remains the final policy mandate.",
                "Trace-centric audit signals are preserved across the lifecycle.",
            ],
        ),
        (
            "Orchestration Sequence",
            diagram_dir / "orchestration-sequence.png",
            summary_b[0] if summary_b else "Deterministic stage progression with bounded retry and fail-safe exits.",
            feature_bullets[:5] if feature_bullets else ["Feature highlights unavailable in source document."],
        ),
        (
            "YAML Pluggability Control Surface",
            diagram_dir / "yaml-pluggability-control-surface.png",
            summary_c[0] if summary_c else "Configuration-first model for runtime agility and controlled rollout.",
            yaml_bullets[:5] if yaml_bullets else ["YAML deep-dive content unavailable in source document."],
        ),
    ]

    missing = [str(path) for _, path, _, _ in slides if not path.exists()]
    if missing:
        raise FileNotFoundError(f"Missing required image files: {missing}")

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(
        prs,
        "Agentic Onboarding Visual Architecture",
        f"{subtitle} | Generated {timestamp}",
    )
    add_agenda_slide(prs)

    for title, image, summary, bullets in slides:
        add_image_slide(prs, title, image, summary, bullets)

    add_text_highlight_slide(prs, "Implementation Anchors", impl_bullets[:6])

    out_name = f"AGENTIC_ONBOARDING_VISUAL_DECK_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    out_path = repo / out_name
    prs.save(out_path)
    print(out_path)


if __name__ == "__main__":
    main()
